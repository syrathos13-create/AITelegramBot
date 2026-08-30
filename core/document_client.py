import io
from docx import Document
from docx.shared import Inches
from openpyxl import Workbook, load_workbook
from pptx import Presentation

def create_docx(title, sections, image_bytes=None):
    doc = Document()
    doc.add_heading(title, level=0)

    if image_bytes:
        doc.add_picture(io.BytesIO(image_bytes), width=Inches(5.5))

    for section in sections:
        if section.get("heading"):
            doc.add_heading(section["heading"], level=1)
        doc.add_paragraph(section.get("text", ""))
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    buf.name = "document.docx"
    return buf

def create_xlsx(title, headers, rows):
    wb = Workbook()
    ws = wb.active
    ws.title = (title or "Sheet1")[:30]
    if headers:
        ws.append(headers)
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    buf.name = "table.xlsx"
    return buf

def _set_slide_background(slide, image_bytes, width, height):
    pic = slide.shapes.add_picture(io.BytesIO(image_bytes), 0, 0, width=width, height=height)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

def create_pptx(title, slides, image_bytes=None):
    prs = Presentation()
    width, height = prs.slide_width, prs.slide_height

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if image_bytes:
        _set_slide_background(title_slide, image_bytes, width, height)

    bullet_layout = prs.slide_layouts[1]
    for s in slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = s.get("title", "")
        bullets = s.get("bullets", [])
        if bullets:
            body = slide.placeholders[1].text_frame
            body.text = bullets[0]
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = b
        if image_bytes:
            _set_slide_background(slide, image_bytes, width, height)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    buf.name = "presentation.pptx"
    return buf

def read_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def read_xlsx(file_bytes):
    wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            lines.append(" | ".join(str(c) if c is not None else "" for c in row))
    return "\n".join(lines)

def read_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        lines.append(text)
    return "\n".join(lines)